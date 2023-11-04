from openpyxl import load_workbook
from pptx import Presentation

# Load Excel file
excel_file_path = 'x.xlsx'
wb = load_workbook(excel_file_path)
sheet = wb.active

# Load PowerPoint template
ppt_template_path = 'y.pptx'
template = Presentation(ppt_template_path)

# Placeholder text in the PowerPoint template
placeholder_text = "{{Name}}"

# Iterate through each row in the Excel sheet
for row in sheet.iter_rows(min_row=2, values_only=True):  # Assuming names start from row 2
    name = row[0]  # Assuming the name is in the first column

    # Duplicate the template
    new_ppt = Presentation(ppt_template_path)

    # Replace placeholder with the name in each slide
    for slide in new_ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                shape.text = shape.text.replace(placeholder_text, name)

    # Save the modified PowerPoint file
    output_path = f'/{name}_certificate.pptx'  # Change the output directory as needed
    new_ppt.save(output_path)
    print('generated for ' + name)
print("Certificates generated successfully!")
